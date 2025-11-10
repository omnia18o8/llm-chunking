""" Document Ingestion Stage 
%pip install pymupdf4llm docx2txt python-pptx """
# ingestion.py

import os
import docx2txt
import pymupdf4llm
from pptx import Presentation

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def load_documents(base_dir):
    folders = [os.path.join(base_dir, f"Module {i}") for i in [1, 2, 4, 5]]
    md_texts = []

    for folder in folders:
        for root, _, files in os.walk(folder):
            for filename in files:
                file_path = os.path.join(root, filename)
                ext = os.path.splitext(file_path)[-1].lower()

                try:
                    if ext == ".pdf":
                        text = pymupdf4llm.to_markdown(file_path)
                    elif ext == ".docx":
                        text = docx2txt.process(file_path)
                    elif ext == ".pptx":
                        text = read_pptx(file_path)
                    else:
                        continue

                    md_texts.append({"file": file_path, "text": text})

                except:
                    continue
    return md_texts
